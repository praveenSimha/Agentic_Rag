from docx import Document
import csv, os
import pptx
import PyPDF2

# === Individual Parsers ===

def parse_pdf(file_path):
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        return "\n".join([page.extract_text() or "" for page in reader.pages])

def parse_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def parse_csv(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        return "\n".join(["\t".join(row) for row in reader])

def parse_pptx(file_path):
    prs = pptx.Presentation(file_path)
    return "\n".join([
        shape.text for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    ])

# === Ingestion Wrapper ===

def ingest(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == 'pdf':
        content = parse_pdf(file_path)
    elif ext == 'docx':
        content = parse_docx(file_path)
    elif ext in ['txt', 'md']:
        content = parse_txt(file_path)
    elif ext == 'csv':
        content = parse_csv(file_path)
    elif ext == 'pptx':
        content = parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")
    return content

# === MCP-Aware Ingestion Handler ===

def handle_ingestion(mcp_message):
    file_path = mcp_message["payload"]["file_path"]
    trace_id = mcp_message["trace_id"]

    content = ingest(file_path)

    # Simple fixed-size chunking
    chunks = [content[i:i+1000] for i in range(0, len(content), 1000)]

    return {
        "sender": "IngestionAgent",
        "receiver": "RetrievalAgent",
        "type": "CHUNK_RESPONSE",
        "trace_id": trace_id,
        "payload": {"chunks": chunks}
    }
